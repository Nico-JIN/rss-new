from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

slides = [
    {
        "type": "title",
        "title": "OpenClaw 数字龙虾概览",
        "subtitle": "是什么｜有什么用｜怎么用｜安全治理\n2026年4月13日"
    },
    {
        "type": "content",
        "title": "OpenClaw 数字龙虾是什么",
        "bullets": [
            "OpenClaw（昵称“数字龙虾”或“养龙虾”）是近期爆火的AI代理平台，将复杂任务封装成可持续运行的虚拟员工。",
            "核心能力：调度大模型 + 工具链 + RPA脚本，形成 24x7 自主循环的“AI员工”以处理信息采集、写作、客服等场景。",
            "形态覆盖桌面端、浏览器的 Agent Phone，以及企业版 ClawPro 套件，可在本地或云端部署并细化权限。"
        ]
    },
    {
        "type": "content",
        "title": "它能带来什么价值",
        "bullets": [
            "效率：在资讯监测、行业简报、代码审查等长周期任务上不间断运行，输出频率可提升 3-10 倍。",
            "创作：协助撰写专利草稿、营销方案、舆情回应模板，但需结合人工校对和法律审查。",
            "运营：串联CRM、工单、表格等系统，自动同步数据、发起审批或提醒，提高跨部门协作透明度。",
            "创新：通过多代理协作，把市场调研、视觉创作、财务测算等步骤拆分并并行执行。"
        ]
    },
    {
        "type": "content",
        "title": "怎么用：典型落地流程",
        "bullets": [
            "1. 账号与权限：申请官方账号或自建实例，配置API秘钥与模型配额，限定可访问的数据域。",
            "2. 工作台搭建：选定任务模版（资讯巡检、合规复核等），补充自有Prompt、工具脚本及触发频率。",
            "3. 数据/工具接入：连接企业知识库、日程、邮箱、协同工具，设置输入输出格式。",
            "4. 监控与回路：通过仪表盘观察token、延时与产出质量，人工抽检或二次模型校验后再推送结果。"
        ]
    },
    {
        "type": "content",
        "title": "安全与合规要点",
        "bullets": [
            "数据分级：涉及客户、专利或财务信息时启用脱敏与最小访问策略，必要时放在内网沙箱运行。",
            "模型访问控制：限制可调用的模型和第三方插件，记录所有调用日志以便审计与费用管理。",
            "Prompt 注入防护：为关键任务设置系统提示词白名单，对外部网页/邮件内容先做过滤再交给代理。",
            "政策遵循：关注监管提醒——例如知识产权部门已提示用“龙虾”撰写专利存在泄密和合规风险。",
            "人机协同：关键决策安排人工复核，把代理输出视为草稿，建立“异常即回滚”的快速止损机制。"
        ]
    },
    {
        "type": "content",
        "title": "推进路径与资源建议",
        "bullets": [
            "试点分层：先在资讯监测、客户FAQ等低风险场景打造MVP，再扩展到研发、法务等高价值链路。",
            "效果量化：为每个“数字龙虾”设置可量化OKR（节省人时、召回率、响应时间），形成投资回报闭环。",
            "安全左移：搭建统一的密钥管控、日志平台和合规检查表，形成标准化接入流程。",
            "赋能团队：培训“驯养师”角色，掌握任务拆解、Prompt设计、故障排查，保障持续运营。"
        ]
    }
]

for slide in slides:
    if slide["type"] == "title":
        layout = prs.slide_layouts[0]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide["title"]
        s.placeholders[1].text = slide["subtitle"]
    else:
        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide["title"]
        body = s.shapes.placeholders[1].text_frame
        body.text = slide["bullets"][0]
        for bullet in slide["bullets"][1:]:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0
        for paragraph in body.paragraphs:
            paragraph.font.size = Pt(20)

output_path = r"C:\Users\76539\.openclaw\skills\rss-news\openclaw_shuzi_longxia_intro.pptx"
prs.save(output_path)
print(output_path)
